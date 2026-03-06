import asyncio
import io
import logging
import pathlib
import re
from typing import Literal

import PyPDF2
from google import genai  # Modern Gemini SDK
from docx import Document
from PIL import Image
from pptx import Presentation

from app.core.config import settings

logger = logging.getLogger(__name__)

CategoryType = Literal["SYLLABUS", "PYQ", "NOTES"]

# --- OCR CONFIGURATION ---
_OCR_PROMPT = "Transcribe all text from this image accurately. Output only the transcribed text."
_MAX_RETRIES = 3
_RETRY_WAIT = 65.0  # Safe margin for free-tier rate limits

class _GeminiRateLimiter:
    """
    Proactive Throttle: Enforces a minimum interval between AI calls.
    Designed to stay safely under the 15 RPM free-tier limit.
    """
    def __init__(self, min_interval: float = 7.0) -> None:
        self._min_interval = min_interval
        self._last_call: float = 0.0
        self._lock = asyncio.Lock()

    async def wait(self) -> None:
        async with self._lock:
            loop = asyncio.get_event_loop()
            now = loop.time()
            gap = self._min_interval - (now - self._last_call)
            if gap > 0:
                print(f"[THROTTLE] Cooling down: {gap:.1f}s remaining...", flush=True)
                await asyncio.sleep(gap)
            self._last_call = loop.time()

# Shared instance: 7s interval = ~8.5 Requests Per Minute (RPM)
_gemini_rate_limiter = _GeminiRateLimiter(min_interval=7.0)

class UniversalExtractor:
    """
    The 'Data Ingestion' layer. Dynamically parses PDF, DOCX, PPTX, and Images.
    Uses Gemini Vision for high-accuracy OCR on handwritten or scanned notes.
    """
    def __init__(self) -> None:
        # Initialising the modern Google GenAI Client
        self.client = genai.Client(
            api_key=settings.GEMINI_API_KEY,
            http_options={"api_version": "v1"},
        )
        self.model_id = "gemini-1.5-flash"

    async def extract_content_from_bytes(self, file_bytes: bytes, filename: str, category: CategoryType) -> str:
        """
        Main entry point for file processing.
        """
        ext = pathlib.Path(filename).suffix.lower()
        buf = io.BytesIO(file_bytes)
        
        print(f"[EXTRACTOR] Processing: {filename} | Category: {category}", flush=True)
        
        try:
            if ext == ".pdf":
                text = await asyncio.to_thread(self._extract_pdf, buf)
            elif ext == ".docx":
                text = await asyncio.to_thread(self._extract_docx, buf)
            elif ext == ".pptx":
                text = await asyncio.to_thread(self._extract_pptx, buf)
            elif ext in {".jpg", ".jpeg", ".png"}:
                # Images require OCR via Gemini Vision
                text = await self._extract_image_ocr_bytes(buf)
            else:
                text = f"[Warning: Unsupported format {ext}]"
        except Exception as e:
            logger.error(f"Extraction failed for {filename}: {e}")
            text = f"[Error parsing file content: {str(e)}]"

        # Structured output for AI context injection
        return (
            f"\n\n--- SOURCE: {filename} ---\n"
            f"TYPE: {category}\n"
            f"{text.strip()}\n"
            f"--- END SOURCE ---\n\n"
        )

    def _extract_pdf(self, buf: io.BytesIO) -> str:
        """Extracts text from standard PDF layers."""
        reader = PyPDF2.PdfReader(buf)
        return "\n\n".join([p.extract_text() or "" for p in reader.pages])

    def _extract_docx(self, buf: io.BytesIO) -> str:
        """Extracts text from MS Word paragraphs."""
        doc = Document(buf)
        return "\n".join([p.text for p in doc.paragraphs])

    def _extract_pptx(self, buf: io.BytesIO) -> str:
        """Extracts text from PowerPoint slide shapes."""
        prs = Presentation(buf)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    async def _extract_image_ocr_bytes(self, buf: io.BytesIO) -> str:
        """
        AI-Powered OCR: Uses Gemini 1.5 Flash to 'read' images.
        Includes built-in retry logic for rate-limit handling.
        """
        def _call_vision_api():
            buf.seek(0)
            img = Image.open(buf)
            # Modern SDK syntax for multimodal generation
            response = self.client.models.generate_content(
                model=self.model_id,
                contents=[_OCR_PROMPT, img],
            )
            return response.text

        for attempt in range(1, _MAX_RETRIES + 1):
            await _gemini_rate_limiter.wait() # Ensure RPM compliance
            try:
                print(f"[OCR] Analysing image structure (Attempt {attempt})...", flush=True)
                result = await asyncio.to_thread(_call_vision_api)
                return result
            except Exception as exc:
                if "429" in str(exc) and attempt < _MAX_RETRIES:
                    print(f"[OCR] Rate limited! Waiting {_RETRY_WAIT}s...", flush=True)
                    await asyncio.sleep(_RETRY_WAIT)
                else:
                    raise